import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import datetime

def create_pptx(news_list: list, output_filename: str = None) -> str:
    """
    將新聞列表生成一文一頁的 PPTX 簡報。
    如果傳入的是 SQLAlchemy Object，需預先轉換為 dict，或直接使用屬性。
    此處假設 news_list 為 dict 的列表。
    """
    prs = Presentation()
    
    # 建立標題頁
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    today_str = datetime.date.today().strftime("%Y-%m-%d")
    title.text = f"AI 自動化新聞簡報 - {today_str}"
    subtitle.text = f"共收錄 {len(news_list)} 篇重點新聞\n生成時間: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    # 建立內容頁 (一文一頁)
    # 版面 1 通常是 Title and Content
    bullet_slide_layout = prs.slide_layouts[1] 
    
    for news in news_list:
        # 兼容 dict 或 SQLAlchemy 模型物件
        if isinstance(news, dict):
            n_title = news.get('title', '無標題')
            n_summary = news.get('summary', '無摘要內容')
            n_source = news.get('source', '未知來源')
            n_category = news.get('category', '未分類')
            n_url = news.get('url', '')
        else:
            n_title = getattr(news, 'title', '無標題')
            n_summary = getattr(news, 'summary', '無摘要內容')
            n_source = getattr(news, 'source', '未知來源')
            n_category = getattr(news, 'category', '未分類')
            n_url = getattr(news, 'url', '')

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        # 設定標題
        title_shape = shapes.title
        title_shape.text = n_title
        
        # 設定摘要內容
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # 清除預設內容
        
        # 簡單分割行來處理條列
        lines = n_summary.split('\n')
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
                
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = line
            p.font.size = Pt(16)
            
        # 加上來源與時間標示 (頁尾)
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(9)
        height = Inches(0.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf_footer = txBox.text_frame
        
        p_footer = tf_footer.add_paragraph()
        p_footer.text = f"分類: {n_category} | 來源: {n_source}\n連結: {n_url}"
        p_footer.font.size = Pt(12)
        p_footer.font.color.rgb = RGBColor(128, 128, 128)
        
    if not output_filename:
        output_filename = f"News_Brief_{today_str}.pptx"
        
    output_path = os.path.join(os.path.dirname(__file__), output_filename)
    prs.save(output_path)
    
    return output_path
